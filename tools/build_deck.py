"""
ON-FARM 제출용 발표 자료(데이터 분석 기획서)를 만든다.

    python tools/build_deck.py

원칙(이 파일이 지키는 것):
  1. 숫자는 코드가 계산한다. 슬라이드에 손으로 적은 수치가 없어야
     발표 직전에 바뀐 값이 조용히 어긋나지 않는다.
  2. 모든 수치 슬라이드 하단에 출처를 박는다. 「데이터 이해·분석 20점」은
     숫자의 크기가 아니라 출처와 한계를 아는지를 본다.
  3. 화면 이미지는 docs/screens/ 의 실제 캡처만 쓴다(tools/capture_screens.mjs).
     목업을 넣으면 시연에서 드러난다.

산출:
  outputs/ON-FARM_천안_기획서.pptx
  outputs/ON-FARM_천안_기획서.pdf   (LibreOffice 있을 때)
  outputs/charts/*.png
"""

from __future__ import annotations

import json
import re
import subprocess
import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Windows 콘솔 기본 코드페이지(cp949)에서 기호가 깨지지 않게 한다
for stream in (sys.stdout, sys.stderr):
    if hasattr(stream, "reconfigure"):
        stream.reconfigure(encoding="utf-8", errors="replace")

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "outputs"
CHARTS = OUT / "charts"
SCREENS = ROOT / "docs" / "screens"
OUT.mkdir(exist_ok=True)
CHARTS.mkdir(exist_ok=True)

# ── 수치는 전부 여기 한 곳에서 온다 ────────────────────────────────────
CHEONAN_FARMS = 9_488          # KOSIS DT_1EA1011 (2024), 상대표준오차 4.5%
CHEONAN_FARM_POP = 21_433
LOCALFOOD_FARMS = 2_256        # 한국경제 2026-05-27
LOCALFOOD_STORES = 12
LOCALFOOD_SALES_EOK = 271
ELDERLY_RATE = 0.700           # 천안 실측 7,773/11,111 (DT_1AG20107, 2020 농림어업총조사)
CHEONAN_MEAN_AGE = 65.1        # 위와 같은 표, 경영주 평균 연령
SMALL_FARM_RATE = 0.767        # 위와 같은 표, 경지 1ha 미만 8,524/11,111
CONTEST_WORKS = 173
CONTEST_BLANK = 6

OUTSIDE = CHEONAN_FARMS - LOCALFOOD_FARMS
OUTSIDE_RATE = OUTSIDE / CHEONAN_FARMS
ELDERLY_OUTSIDE = round(OUTSIDE * ELDERLY_RATE)

META = json.loads((ROOT / "models" / "metadata.json").read_text(encoding="utf-8"))
PER_ITEM: dict = META.get("per_item") or {}
POLICY_CAP = 0.85              # src/ai/policy.ts POLICY_MAX_WITHOUT_FIELD_EVAL


def test_count() -> int:
    """
    테스트 개수를 소스에서 '세지' 않고 실제 실행 결과에서 읽는다.

    파일을 세면 중첩 describe 나 표기 차이 때문에 조용히 어긋난다(실제로 150 vs 156 로 어긋났다).
    겸사겸사, 하나라도 실패하면 발표 자료가 만들어지지 않는다 —
    깨진 코드로 만든 슬라이드를 들고 나가는 일이 없게 한다.
    """
    proc = subprocess.run(
        ["npm.cmd" if sys.platform == "win32" else "npm", "test"],
        cwd=ROOT, capture_output=True, text=True, encoding="utf-8",
        errors="replace", timeout=900,
    )
    out = proc.stdout + proc.stderr
    got = {
        key: int(m.group(1))
        for key in ("tests", "pass", "fail")
        if (m := re.search(rf"^\D*\b{key} (\d+)$", out, re.MULTILINE))
    }
    if "tests" not in got:
        raise SystemExit("테스트 결과를 읽지 못했습니다. `npm test` 를 직접 실행해 확인하세요.")
    if got.get("fail"):
        raise SystemExit(f"테스트 {got['fail']}개가 실패 중입니다. 고친 뒤에 자료를 만드세요.")
    return got["tests"]


TESTS = test_count()


def check_docs_agree(tests: int) -> None:
    """
    문서에 박아 둔 테스트 개수가 실제와 어긋나면 자료를 만들지 않는다.

    실제로 세 문서가 138·148·136 으로 서로 달랐다. 발표장에서 심사위원이
    "테스트 몇 개라고요?" 하고 되물으면 그 순간 다른 숫자도 의심받는다.
    50 이상인 값만 '전체 개수 주장'으로 보고 검사한다(부분 개수 서술은 통과).
    """
    stale = []
    targets = [ROOT / "README.md", ROOT / "DEVELOPMENT_STATUS.md", *(ROOT / "docs").glob("*.md")]
    for f in targets:
        if not f.exists():
            continue
        for m in re.finditer(r"테스트 (\d+)개", f.read_text(encoding="utf-8")):
            n = int(m.group(1))
            if n >= 50 and n != tests:
                stale.append(f"{f.relative_to(ROOT)}: '테스트 {n}개'")
    if stale:
        raise SystemExit(
            f"실제 테스트는 {tests}개인데 문서가 다르게 적고 있습니다:\n  "
            + "\n  ".join(stale)
        )


check_docs_agree(TESTS)

# ── 색: 팜보이스 덱(짙은 녹색+라임)과 겹치지 않게 앱 자체 색을 쓴다 ──
CREAM = RGBColor(0xFD, 0xF8, 0xF0)
INK = RGBColor(0x23, 0x20, 0x1C)
MUTED = RGBColor(0x7A, 0x72, 0x66)
ORANGE = RGBColor(0xDD, 0x4B, 0x14)
GREEN = RGBColor(0x1F, 0x7A, 0x4D)
LINE = RGBColor(0xE6, 0xDD, 0xCC)
RED = RGBColor(0xB4, 0x25, 0x25)

HEX = lambda c: f"#{c[0]:02X}{c[1]:02X}{c[2]:02X}"

for name in ("Malgun Gothic", "맑은 고딕", "NanumGothic"):
    if any(f.name == name for f in font_manager.fontManager.ttflist):
        plt.rcParams["font.family"] = name
        break
plt.rcParams["axes.unicode_minus"] = False
FONT = "맑은 고딕"

W, H = Inches(13.333), Inches(7.5)


# ── 슬라이드 조립 도우미 ──────────────────────────────────────────────
def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs, bg=CREAM):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def text(slide, x, y, w, h, runs, *, size=18, color=INK, bold=False,
         align=PP_ALIGN.LEFT, spacing=1.25, anchor=MSO_ANCHOR.TOP):
    """runs: 문자열 또는 (문자열, {size,color,bold}) 목록. 문단은 리스트로 나눈다."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        chunks = para if isinstance(para, list) else [para]
        for chunk in chunks:
            body, opt = chunk if isinstance(chunk, tuple) else (chunk, {})
            r = p.add_run()
            r.text = body
            r.font.name = FONT
            r.font.size = Pt(opt.get("size", size))
            r.font.bold = opt.get("bold", bold)
            r.font.color.rgb = opt.get("color", color)
    return box


def box(slide, x, y, w, h, *, fill=None, line=LINE, radius=True, width=Pt(1.25)):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    if radius:
        shp.adjustments[0] = 0.06
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = width
    shp.shadow.inherit = False
    return shp


def head(slide, kicker, title, n):
    """모든 본문 슬라이드의 머리 — 위치를 고정해 넘길 때 흔들리지 않게 한다."""
    text(slide, Inches(0.7), Inches(0.42), Inches(11.9), Inches(0.32),
         kicker, size=13, color=ORANGE, bold=True)
    text(slide, Inches(0.7), Inches(0.75), Inches(11.9), Inches(0.75),
         title, size=31, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.52), Inches(11.93), Pt(1.6))
    line.fill.solid(); line.fill.fore_color.rgb = LINE; line.line.fill.background()
    line.shadow.inherit = False
    text(slide, Inches(12.1), Inches(6.92), Inches(0.6), Inches(0.3),
         str(n), size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def source(slide, note):
    text(slide, Inches(0.7), Inches(6.88), Inches(11.2), Inches(0.4),
         note, size=10.5, color=MUTED, spacing=1.1)


def picture(slide, name, x, y, *, height=None, width=None, border=True):
    """실제 캡처만 넣는다. 없으면 빌드를 실패시켜 목업이 섞이지 않게 한다."""
    path = SCREENS / f"{name}.png"
    if not path.exists():
        raise SystemExit(
            f"화면 캡처가 없습니다: {path}\n"
            "먼저 서버를 켜고 tools/capture_screens.mjs 를 실행하세요."
        )
    kw = {}
    if height: kw["height"] = height
    if width: kw["width"] = width
    pic = slide.shapes.add_picture(str(path), x, y, **kw)
    if border:
        edge = box(slide, pic.left, pic.top, pic.width, pic.height,
                   fill=None, line=LINE, width=Pt(1))
        edge.shadow.inherit = False
    return pic


# ── 차트 ──────────────────────────────────────────────────────────────
def chart_gap() -> Path:
    """천안 농가 9,488 중 로컬푸드 밖 7,232 — 제안의 근거 한 장."""
    fig, ax = plt.subplots(figsize=(5.6, 4.4), dpi=200)
    fig.patch.set_facecolor(HEX(CREAM))
    ax.set_facecolor(HEX(CREAM))
    wedges, _ = ax.pie(
        [OUTSIDE, LOCALFOOD_FARMS],
        startangle=90, counterclock=False,
        colors=[HEX(ORANGE), HEX(GREEN)],
        wedgeprops=dict(width=0.34, edgecolor=HEX(CREAM), linewidth=3),
    )
    ax.text(0, 0.12, f"{OUTSIDE_RATE:.1%}", ha="center", va="center",
            fontsize=40, fontweight="bold", color=HEX(INK))
    ax.text(0, -0.22, "로컬푸드 체계 밖", ha="center", va="center",
            fontsize=13, color=HEX(MUTED))
    ax.text(0, -0.42, f"{OUTSIDE:,}가구", ha="center", va="center",
            fontsize=15, fontweight="bold", color=HEX(ORANGE))
    ax.set(aspect="equal")
    fig.tight_layout(pad=0.2)
    p = CHARTS / "gap.png"
    fig.savefig(p, facecolor=fig.get_facecolor()); plt.close(fig)
    return p


def chart_grade() -> Path:
    """등급 정확도를 '중량만 쓴 기준선'과 비교한다 — 공정 대조군."""
    items = list(PER_ITEM.keys())
    acc = [PER_ITEM[i].get("grade_object_acc", 0) for i in items]
    base = [PER_ITEM[i].get("weight_only_baseline", 0) for i in items]
    usable = [PER_ITEM[i].get("grade_usable", False) for i in items]

    # 슬라이드에서 아래 주석 상자를 덮지 않도록 납작하게 만든다
    fig, ax = plt.subplots(figsize=(10.6, 3.0), dpi=200)
    fig.patch.set_facecolor(HEX(CREAM)); ax.set_facecolor(HEX(CREAM))
    x = range(len(items)); w = 0.36
    ax.bar([i - w / 2 for i in x], acc, w, label="AI 등급 정확도",
           color=[HEX(GREEN) if u else HEX(RED) for u in usable])
    ax.bar([i + w / 2 for i in x], base, w, label="중량만 쓴 기준선",
           color="#C9BFAE")
    for i, (a, b, u) in enumerate(zip(acc, base, usable)):
        ax.text(i - w / 2, a + 0.03, f"{a:.0%}", ha="center", fontsize=11,
                fontweight="bold", color=HEX(GREEN) if u else HEX(RED))
        ax.text(i + w / 2, b + 0.03, f"{b:.0%}", ha="center", fontsize=10, color=HEX(MUTED))
    # '표시 차단'을 축 아래 별도 글자로 두면 슬라이드에서 잘린다 — 눈금 이름에 붙인다
    ax.set_xticks(list(x))
    ax.set_xticklabels(
        [f"{n}  ·  표시 차단" if not u else n for n, u in zip(items, usable)], fontsize=13
    )
    for label, u in zip(ax.get_xticklabels(), usable):
        if not u:
            label.set_color(HEX(RED)); label.set_fontweight("bold")
    ax.set_ylim(0, 1.12); ax.set_yticks([0, 0.5, 1.0])
    ax.set_yticklabels(["0%", "50%", "100%"], fontsize=10, color=HEX(MUTED))
    # 범례를 그림 안에 두면 100% 막대의 값 표시와 겹친다 — 축 위로 뺀다
    ax.legend(fontsize=10.5, frameon=False, loc="lower left",
              bbox_to_anchor=(0, 1.0), ncol=2)
    for s in ("top", "right"): ax.spines[s].set_visible(False)
    ax.spines["left"].set_color(HEX(LINE)); ax.spines["bottom"].set_color(HEX(LINE))
    ax.tick_params(length=0)
    fig.tight_layout(pad=0.3)
    p = CHARTS / "grade.png"
    fig.savefig(p, facecolor=fig.get_facecolor()); plt.close(fig)
    return p


# ── 슬라이드 ──────────────────────────────────────────────────────────
def build() -> Path:
    prs = new_deck()
    gap_png, grade_png = chart_gap(), chart_grade()

    # 1 표지 ────────────────────────────────────────────────────────────
    s = blank(prs, INK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.05), Inches(0.11), Inches(1.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background(); bar.shadow.inherit = False
    text(s, Inches(0.9), Inches(1.25), Inches(11), Inches(0.4),
         "2026년 천안시 AI·데이터 기반 정책 아이디어 경진대회 · 과제4 스마트 농업 & 축산 생산성 향상",
         size=14, color=RGBColor(0xC9, 0xBF, 0xAE))
    text(s, Inches(1.25), Inches(2.0), Inches(11), Inches(1.7),
         [[("ON-FARM", {"size": 58, "bold": True, "color": CREAM})],
          [("사진 한 장으로 끝나는 고령농 출하 도구", {"size": 26, "color": RGBColor(0xE8, 0xE0, 0xD2)})]],
         spacing=1.15)
    text(s, Inches(0.9), Inches(4.15), Inches(11.2), Inches(1.4),
         [[("기른 것을 못 파는 손실을 줄여 ", {"size": 21, "color": RGBColor(0xE8, 0xE0, 0xD2)}),
           ("이미 심은 농산물의 실현 생산성", {"size": 21, "bold": True, "color": ORANGE}),
           ("을 높입니다.", {"size": 21, "color": RGBColor(0xE8, 0xE0, 0xD2)})]],
         spacing=1.3)
    for i, (k, v) in enumerate([
        ("천안 농가 중 로컬푸드 체계 밖", f"{OUTSIDE_RATE:.1%}"),
        ("농가가 입력하는 항목", "0개"),
        ("작동하는 MVP · 자동 테스트", f"{TESTS}개"),
    ]):
        x = Inches(0.9 + i * 3.95)
        text(s, x, Inches(5.55), Inches(3.7), Inches(0.34), k, size=12,
             color=RGBColor(0x9A, 0x92, 0x85))
        text(s, x, Inches(5.86), Inches(3.7), Inches(0.6), v, size=27, bold=True, color=CREAM)
    text(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4),
         "제출: 2026년 8월 · 본 자료의 모든 수치는 공개 데이터에서 재계산하며 산출 스크립트를 함께 제출합니다.",
         size=11, color=RGBColor(0x7A, 0x72, 0x66))

    # 2 실현 생산성 프레임 ───────────────────────────────────────────────
    s = blank(prs); head(s, "왜 이것이 과제4의 '생산성' 인가", "생산성은 밭에서 끝나지 않는다", 2)
    text(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(0.6),
         [[("스마트팜·품종개량은 ", {}), ("수확량", {"bold": True}),
           ("까지를 늘립니다. 그런데 천안의 고령 소농에게 손실이 나는 곳은 그다음 칸입니다.", {})]],
         size=17.5, color=MUTED)
    stages = [("잠재 생산량", "땅과 품종"), ("수확량", "재배 기술"), ("실제 판매량", "출하 능력")]
    for i, (t, sub) in enumerate(stages):
        x = Inches(0.85 + i * 4.1)
        last = i == 2
        box(s, x, Inches(2.75), Inches(3.35), Inches(1.55),
            fill=CREAM if not last else RGBColor(0xFD, 0xEC, 0xE3),
            line=ORANGE if last else LINE, width=Pt(2.5 if last else 1.25))
        text(s, x, Inches(3.02), Inches(3.35), Inches(0.5), t, size=21, bold=True,
             align=PP_ALIGN.CENTER, color=ORANGE if last else INK)
        text(s, x, Inches(3.55), Inches(3.35), Inches(0.4), sub, size=13,
             align=PP_ALIGN.CENTER, color=MUTED)
        if i < 2:
            text(s, Inches(4.15 + i * 4.1), Inches(3.05), Inches(0.75), Inches(0.5),
                 "→", size=26, align=PP_ALIGN.CENTER, color=MUTED)
    text(s, Inches(9.05), Inches(4.4), Inches(3.35), Inches(0.45),
         "▲ 여기가 비어 있다", size=15, bold=True, align=PP_ALIGN.CENTER, color=ORANGE)
    box(s, Inches(0.85), Inches(5.15), Inches(11.6), Inches(1.35), fill=RGBColor(0xF5, 0xEF, 0xE4), line=None)
    text(s, Inches(1.2), Inches(5.42), Inches(11), Inches(0.9),
         [[("밭에서 아무리 잘 길러도 ", {}), ("팔지 못하면 생산성은 0", {"bold": True, "color": ORANGE}),
           ("입니다.", {})],
          [("ON-FARM 은 왼쪽 화살표가 아니라 ", {"color": MUTED}),
           ("오른쪽 화살표", {"bold": True, "color": INK}),
           ("를 굵게 만듭니다 — 이것이 실현 생산성입니다.", {"color": MUTED})]],
         size=17, spacing=1.35)
    source(s, "과제4 「스마트 농업 & 축산 생산성 향상」의 생산성을 재배 단계가 아니라 실현 단계로 정의한다.")

    # 3 천안 근거 ────────────────────────────────────────────────────────
    s = blank(prs); head(s, "천안 데이터", f"천안 농가 4곳 중 3곳은 로컬푸드 체계 밖에 있다", 3)
    s.shapes.add_picture(str(gap_png), Inches(0.75), Inches(1.95), height=Inches(4.4))
    rows = [
        ("천안시 전체 농가", f"{CHEONAN_FARMS:,}가구", "KOSIS 농림어업조사 2024 (상대표준오차 4.5%)"),
        ("로컬푸드 참여 농가", f"{LOCALFOOD_FARMS:,}가구", f"직매장 {LOCALFOOD_STORES}곳 · 2026년 매출 {LOCALFOOD_SALES_EOK}억 원"),
        ("체계 밖 농가", f"{OUTSIDE:,}가구",
         f"그들은 평균 {CHEONAN_MEAN_AGE}세가 경영하는 소농이다 "
         f"(60대 이상 {ELDERLY_RATE:.0%} · 경지 1ha 미만 {SMALL_FARM_RATE:.0%})"),
    ]
    for i, (k, v, note) in enumerate(rows):
        y = Inches(2.1 + i * 1.28)
        text(s, Inches(6.5), y, Inches(3.1), Inches(0.4), k, size=15, color=MUTED)
        text(s, Inches(6.5), y + Inches(0.36), Inches(3.1), Inches(0.55), v, size=26, bold=True,
             color=ORANGE if i == 2 else INK)
        text(s, Inches(9.5), y + Inches(0.12), Inches(3.1), Inches(0.9), note, size=11.5,
             color=MUTED, spacing=1.2)
    box(s, Inches(6.5), Inches(5.95), Inches(6.0), Inches(0.78), fill=RGBColor(0xF5, 0xEF, 0xE4), line=None)
    text(s, Inches(6.75), Inches(6.08), Inches(5.6), Inches(0.55),
         "조사 주체·기준이 달라 정밀한 차집합이 아닙니다. 농가 수는 2024년, "
         "연령·경지 비율은 2020년 총조사입니다.",
         size=12.5, color=INK, spacing=1.2)
    source(s, "출처: KOSIS DT_1EA1011(2024) 농가 · DT_1AG20107(2020 농림어업총조사) 연령·경지 · 로컬푸드 현황 보도(2026-05-27) · 산출 tools/cheonan_gap.py·cheonan_age.py")

    # 4 등록 장벽 ────────────────────────────────────────────────────────
    s = blank(prs); head(s, "문제 정의", "판로가 없는 게 아니라, 등록을 못 한다", 4)
    # 오른쪽 강조 상자 아래로 글이 흘러가지 않게 폭을 상자 앞까지만 준다
    text(s, Inches(0.7), Inches(1.82), Inches(6.6), Inches(0.5),
         "직매장도 기존 직거래 앱도, 농가가 다음을 '입력할 수 있다'고 전제합니다.",
         size=16, color=MUTED)
    olds = ["상품명", "카테고리", "중량·규격", "가격 책정", "상세 설명", "사진 편집"]
    for i, t in enumerate(olds):
        x = Inches(0.85 + (i % 3) * 2.05); y = Inches(2.6 + (i // 3) * 0.86)
        box(s, x, y, Inches(1.9), Inches(0.68), fill=RGBColor(0xF2, 0xEC, 0xE2), line=None)
        text(s, x, y + Inches(0.14), Inches(1.9), Inches(0.4), t, size=14,
             align=PP_ALIGN.CENTER, color=RGBColor(0x8A, 0x80, 0x72))
    text(s, Inches(0.85), Inches(4.42), Inches(5.9), Inches(0.5),
         "고령 소농에게는 이 여섯 칸이 곧 장벽입니다.", size=15, bold=True, color=RED)
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.35), Inches(3.0), Inches(0.85), Inches(0.6))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = ORANGE; arrow.line.fill.background(); arrow.shadow.inherit = False
    box(s, Inches(7.5), Inches(2.2), Inches(4.95), Inches(2.8), fill=RGBColor(0xFD, 0xEC, 0xE3), line=ORANGE, width=Pt(2.5))
    text(s, Inches(7.85), Inches(2.42), Inches(4.4), Inches(0.45), "ON-FARM 에서 농가가 하는 일",
         size=14, bold=True, color=ORANGE)
    for i, t in enumerate(["① 사진 찍기", "② 번호 고르기 (1·2·3)", "③ 수량 확인"]):
        text(s, Inches(7.85), Inches(2.95) + Inches(i * 0.52), Inches(4.4), Inches(0.45),
             t, size=19, bold=True)
    text(s, Inches(7.85), Inches(4.53), Inches(4.4), Inches(0.4),
         "그 외 입력 항목 0개", size=14, color=MUTED)
    box(s, Inches(0.85), Inches(5.3), Inches(11.6), Inches(1.2), fill=RGBColor(0xF5, 0xEF, 0xE4), line=None)
    text(s, Inches(1.2), Inches(5.52), Inches(11), Inches(0.8),
         [[("가격은 AI 가 아니라 ", {"color": MUTED}), ("운영자가 미리 정한 표준 가격", {"bold": True}),
           (" 이고, 등급 확정과 포장은 거점이 합니다. 농가는 ", {"color": MUTED}),
           ("가격을 정하지도, 글을 쓰지도 않습니다.", {"bold": True})]],
         size=16, spacing=1.3)
    source(s, "기존 로컬푸드 직매장은 농가가 직접 상품을 갖다 놓고 가격표를 붙이고 재고를 관리하는 구조다 — 대체가 아니라 보완이다.")

    # 5 해법 흐름 ────────────────────────────────────────────────────────
    s = blank(prs); head(s, "해법", "사진 한 장에서 판매 시작까지", 5)
    steps = [
        ("farmer-02-photo", "① 사진 찍기", "카메라 버튼 하나"),
        ("farmer-03-candidates", "② 번호 고르기", "AI 는 후보만 제시"),
        ("farmer-04-quantity", "③ 수량 확인", "가격은 표준 정찰가"),
        ("farmer-06-done", "④ 판매 시작", "거점에 갖다 놓으면 끝"),
    ]
    for i, (img, title, sub) in enumerate(steps):
        x = Inches(0.72 + i * 3.12)
        picture(s, img, x + Inches(0.42), Inches(1.95), height=Inches(3.55))
        text(s, x, Inches(5.65), Inches(2.95), Inches(0.4), title, size=17, bold=True,
             align=PP_ALIGN.CENTER, color=ORANGE)
        text(s, x, Inches(6.05), Inches(2.95), Inches(0.4), sub, size=12.5,
             align=PP_ALIGN.CENTER, color=MUTED)
    source(s, "실제 실행 중인 서버 화면 캡처(재현: npm start → tools/capture_screens.mjs). 목업이 아닙니다.")

    # 6 핵심 화면 ────────────────────────────────────────────────────────
    s = blank(prs); head(s, "핵심 설계", "AI 는 판정하지 않는다. 물어본다", 6)
    picture(s, "farmer-03-candidates", Inches(1.15), Inches(1.85), height=Inches(4.85))
    text(s, Inches(4.7), Inches(2.0), Inches(8.0), Inches(0.9),
         [[("\"이 중에 어느 것인가요?\"", {"size": 27, "bold": True})]], spacing=1.2)
    body = [
        ("확정 대신 질문", "AI 가 100% 확신해도 화면은 항상 번호를 고르게 합니다. 판정 주체는 사람입니다."),
        ("실패해도 진행", "후보에 없으면 '여기 없어요' 로 전체 품목에서 직접 고릅니다. AI 가 틀려도 판매는 멈추지 않습니다."),
        ("눈이 아니라 귀로도", "화면을 읽어주고, 번호는 '일번·이번·삼번' 으로 말합니다. 수량은 말로도 입력합니다."),
        ("입력창이 없다", "글자를 치는 칸이 한 곳도 없습니다. 누르는 것과 고르는 것만 있습니다."),
    ]
    for i, (k, v) in enumerate(body):
        y = Inches(3.0 + i * 0.92)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.7), y + Inches(0.12), Inches(0.13), Inches(0.13))
        dot.fill.solid(); dot.fill.fore_color.rgb = ORANGE; dot.line.fill.background(); dot.shadow.inherit = False
        text(s, Inches(5.0), y, Inches(7.6), Inches(0.85),
             [[(k + "  ", {"size": 16.5, "bold": True}), (v, {"size": 14, "color": MUTED})]],
             spacing=1.3)
    source(s, "화면의 후보는 실제 모델 출력이며, 상위 3개 중 판매 가능한 품목만 제시한다(src/ai/pipeline.ts).")

    # 7 안전 정책 ────────────────────────────────────────────────────────
    s = blank(prs); head(s, "책임 경계", "AI 가 하지 않는 일을 코드가 강제한다", 7)
    cards = [
        ("가격", "AI 가 정하지 않는다", "운영자가 미리 정한 표준 가격만 보여준다. 화면에도 그렇게 적혀 있다."),
        ("등급", "확정이 아닌 참고값", "확정은 거점의 실물 검수다. 근거가 없는 품목은 아예 표시하지 않는다."),
        ("신뢰도", f"{POLICY_CAP:.0%} 를 넘지 못한다", "실환경 평가 전에는 상한이 걸린다. 어떤 판정 엔진을 붙여도 정책이 먼저 적용된다."),
    ]
    for i, (tag, title, desc) in enumerate(cards):
        x = Inches(0.8 + i * 3.98)
        box(s, x, Inches(1.95), Inches(3.7), Inches(2.15), fill=CREAM, line=LINE)
        text(s, x + Inches(0.3), Inches(2.18), Inches(3.1), Inches(0.35), tag, size=12.5,
             bold=True, color=ORANGE)
        text(s, x + Inches(0.3), Inches(2.52), Inches(3.15), Inches(0.6), title, size=18, bold=True)
        text(s, x + Inches(0.3), Inches(3.14), Inches(3.15), Inches(0.85), desc, size=12.5,
             color=MUTED, spacing=1.25)
    picture(s, "farmer-04-quantity", Inches(0.8), Inches(4.35), height=Inches(2.35))
    text(s, Inches(4.05), Inches(4.5), Inches(8.4), Inches(2.1),
         [[("실제 화면에 그대로 적혀 있습니다.", {"size": 16, "bold": True})],
          [("\"이 가격은 운영자가 미리 정해 둔 표준 가격입니다. AI 가 정한 값이 아닙니다.\"",
            {"size": 15.5, "color": ORANGE, "bold": True})],
          [("문서에만 쓴 약속은 시연에서 깨집니다. 그래서 이 경계는 화면 문구와 서버 코드 양쪽에 있고, "
            "자동 테스트가 문구가 사라지는 순간 실패합니다.", {"size": 13.5, "color": MUTED})]],
         spacing=1.4)
    source(s, "src/ai/policy.ts 가 모든 판정 엔진의 출력에 상한과 품목별 등급 게이트를 적용한다. "
              "정책을 우회하면 테스트가 깨지도록 고정했다.")

    # 8 데이터 이해 ① ────────────────────────────────────────────────────
    s = blank(prs); head(s, "데이터 이해 ①", "무엇으로 배웠고, 무엇을 아직 모르는가", 8)
    left = [
        ("학습 데이터", "AI Hub 「농산물 품질(QC) 이미지」 · 5개 품목(사과·배·감귤·감자·양파)"),
        ("평가 단위", "이미지가 아니라 개체(group) 단위. 같은 물체를 여러 장 찍은 데이터라 이미지 단위로 나누면 성능이 부풀려진다."),
        ("품목 인식", "개체 단위 100% — 다만 오른쪽 한계와 함께 읽어야 하는 숫자다."),
    ]
    right = [
        ("스튜디오 촬영", "흰 배경·균일 조명에서 찍은 사진으로 배웠다. 밭에서 찍은 폰 사진은 다르다."),
        ("실환경 미평가", "metadata 의 field_evaluated 가 false 다. 이 값이 true 가 되는 것이 PoC 의 성공 정의다."),
        ("그래서 상한", f"실환경 평가 전에는 신뢰도가 {POLICY_CAP:.0%} 를 넘지 못하게 코드가 막는다."),
    ]
    for col, (items, title, color) in enumerate([(left, "확인된 것", GREEN), (right, "아직 모르는 것", RED)]):
        x = Inches(0.8 + col * 6.05)
        box(s, x, Inches(1.9), Inches(5.75), Inches(4.55), fill=CREAM, line=LINE)
        text(s, x + Inches(0.35), Inches(2.12), Inches(5.0), Inches(0.4), title, size=15, bold=True, color=color)
        for i, (k, v) in enumerate(items):
            y = Inches(2.68 + i * 1.28)
            text(s, x + Inches(0.35), y, Inches(5.05), Inches(1.2),
                 [[(k, {"size": 15, "bold": True})], [(v, {"size": 12.5, "color": MUTED})]],
                 spacing=1.3)
    box(s, Inches(0.8), Inches(6.55), Inches(11.85), Inches(0.0), fill=None, line=None)
    source(s, "숫자를 크게 말하지 않고 한계를 먼저 말한다. 스튜디오 사진으로 잰 정확도를 현장 성능으로 옮겨 적는 순간, "
              "그 제안은 현장에서 무너진다.")

    # 9 데이터 이해 ② ────────────────────────────────────────────────────
    s = blank(prs); head(s, "데이터 이해 ②", "등급은 '중량만 써도 맞히는' 기준선과 비교했다", 9)
    text(s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(0.55),
         "농산물 등급은 상당 부분 크기 등급입니다. 그래서 AI 를 자랑하기 전에, 무게만 쓰는 단순 규칙과 비교했습니다.",
         size=15.5, color=MUTED)
    s.shapes.add_picture(str(grade_png), Inches(0.9), Inches(2.3), width=Inches(11.5))
    onion = PER_ITEM.get("양파", {})
    box(s, Inches(0.9), Inches(5.72), Inches(11.5), Inches(0.95), fill=RGBColor(0xFB, 0xEA, 0xEA), line=None)
    text(s, Inches(1.25), Inches(5.9), Inches(11.0), Inches(0.6),
         [[("양파는 AI 가 기준선보다 못했습니다", {"size": 15.5, "bold": True, "color": RED}),
           (f"  ({onion.get('grade_object_acc', 0):.1%} vs {onion.get('weight_only_baseline', 0):.1%}). "
            "그래서 양파 등급은 화면에 표시되지 않습니다 — 발표에서 숨기지 않고, 코드가 자동으로 막습니다.",
            {"size": 14, "color": INK})]],
         spacing=1.25)
    source(s, "기준선은 학습 분할에서만 임계값을 맞추고 검증 분할에서 한 번만 평가했다(models/metadata.json). "
              "품목별 게이트는 src/ai/policy.ts 의 gradeAllowedFor 가 강제한다.")

    # 10 거점 결합 ───────────────────────────────────────────────────────
    s = blank(prs); head(s, "운영 구조", "기존 로컬푸드 거점을 그대로 쓴다", 10)
    nodes = [
        ("농가", "사진 · 수량", "집·밭에서 3분", GREEN),
        ("로컬푸드 거점", "검수 · 등급 확정 · 포장", "기존 시설·인력", ORANGE),
        ("소비자", "천안 밖 수요", "전국에서 주문", INK),
    ]
    for i, (t, role, sub, c) in enumerate(nodes):
        x = Inches(0.85 + i * 4.35)
        box(s, x, Inches(2.2), Inches(3.55), Inches(2.0), fill=CREAM, line=c, width=Pt(2.2))
        text(s, x, Inches(2.5), Inches(3.55), Inches(0.5), t, size=22, bold=True,
             align=PP_ALIGN.CENTER, color=c)
        text(s, x, Inches(3.05), Inches(3.55), Inches(0.45), role, size=14, align=PP_ALIGN.CENTER)
        text(s, x, Inches(3.5), Inches(3.55), Inches(0.4), sub, size=12, align=PP_ALIGN.CENTER, color=MUTED)
        if i < 2:
            a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5 + i * 4.35), Inches(2.95), Inches(0.6), Inches(0.45))
            a.fill.solid(); a.fill.fore_color.rgb = LINE; a.line.fill.background(); a.shadow.inherit = False
    picture(s, "hub-01-inspect", Inches(0.85), Inches(4.45), height=Inches(2.2))
    text(s, Inches(5.3), Inches(4.5), Inches(7.3), Inches(2.1),
         [[("새 시설을 짓지 않습니다.", {"size": 17, "bold": True})],
          [("직매장 12곳은 경쟁자가 아니라 인프라입니다. ON-FARM 은 직매장에 갈 수 없는 농가와, "
            "직매장이 닿지 않는 지역 밖 수요를 맡습니다.", {"size": 14, "color": MUTED})],
          [("천안 농산물을 천안 밖으로 파는 구조라, 매출은 천안으로 들어옵니다.",
            {"size": 14.5, "bold": True, "color": ORANGE})]],
         spacing=1.4)
    source(s, "거점 화면도 실제 구현이며, 검수 전 상품은 소비자 화면에 노출되지 않는다.")

    # 11 실행 계획 ───────────────────────────────────────────────────────
    s = blank(prs); head(s, "실행 계획", "성환 1거점 · 5농가 · 8주부터 시작한다", 11)
    phases = [
        ("Phase 1 · 8주", "성환 1거점 · 농가 5호", ["실사용 등록과 소요시간 측정",
                                              "등록 시 실제 폰 사진 수집 → 실환경 평가셋",
                                              "검수·출하·정산 1주기 완주"], ORANGE),
        ("Phase 2 · 12주", "3거점 · 농가 30호", ["확인할 것은 성능이 아니라 운영 부하",
                                              "거점 1곳의 하루 처리 한계",
                                              "인접 농가 묶음 배송 성립 여부"], MUTED),
        ("Phase 3", "전 거점", ["Phase 2 에서 손익분기가 나온 뒤 판단",
                             "지금 확산 규모를 말하지 않는다",
                             "근거 없는 숫자는 적지 않는다"], MUTED),
    ]
    for i, (title, scope, items, c) in enumerate(phases):
        x = Inches(0.8 + i * 3.98)
        box(s, x, Inches(1.95), Inches(3.7), Inches(3.5), fill=CREAM,
            line=c if i == 0 else LINE, width=Pt(2.5 if i == 0 else 1.25))
        text(s, x + Inches(0.3), Inches(2.18), Inches(3.1), Inches(0.4), title, size=14, bold=True, color=c)
        text(s, x + Inches(0.3), Inches(2.58), Inches(3.15), Inches(0.5), scope, size=17.5, bold=True)
        for j, it in enumerate(items):
            text(s, x + Inches(0.3), Inches(3.2) + Inches(j * 0.66), Inches(3.15), Inches(0.6),
                 "· " + it, size=12.5, color=MUTED, spacing=1.25)
    box(s, Inches(0.8), Inches(5.65), Inches(11.85), Inches(1.0), fill=RGBColor(0xFD, 0xEC, 0xE3), line=None)
    text(s, Inches(1.15), Inches(5.85), Inches(11.2), Inches(0.6),
         [[("각 단계 끝에 Go/No-Go 를 둡니다. ", {"size": 16, "color": INK}),
           ("기준에 하나라도 미달하면 다음 단계로 가지 않습니다.", {"size": 16, "bold": True, "color": ORANGE})]],
         spacing=1.25)
    source(s, "성환은 배 주산지이며 로컬푸드 거점이 이미 있다. 새 시설을 짓지 않는 것이 이 계획의 전제다.")

    # 12 Go/No-Go ───────────────────────────────────────────────────────
    s = blank(prs); head(s, "중단 기준", "언제 멈출지를 먼저 정한다", 12)
    criteria = [
        ("등록 완주율", "80% 이상", "시작한 농가가 끝까지 등록했는가"),
        ("1건 등록 소요시간", "3분 이내", "중앙값 기준"),
        ("폰 사진 품목 정확도", "85% 이상", "개체 단위, 실제 현장 사진"),
        ("등록 후 판매 성사율", "50% 이상", "올렸는데 안 팔리면 의미가 없다"),
        ("거점 추가 부담", "건당 10분", "초과하면 협조가 끊긴다"),
        ("농가 재사용 의향", "5명 중 4명", "다음에도 쓰겠는가"),
    ]
    for i, (k, target, why) in enumerate(criteria):
        x = Inches(0.8 + (i % 2) * 6.05); y = Inches(1.95 + (i // 2) * 1.42)
        box(s, x, y, Inches(5.75), Inches(1.18), fill=CREAM, line=LINE)
        text(s, x + Inches(0.32), y + Inches(0.16), Inches(3.4), Inches(0.42), k, size=16, bold=True)
        text(s, x + Inches(0.32), y + Inches(0.62), Inches(3.4), Inches(0.4), why, size=12, color=MUTED)
        text(s, x + Inches(3.6), y + Inches(0.32), Inches(1.85), Inches(0.5), target, size=19,
             bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)
    text(s, Inches(0.8), Inches(6.35), Inches(11.9), Inches(0.45),
         "이 문장을 발표에서 그대로 말합니다 — 어느 하나라도 미달하면 확산하지 않습니다.",
         size=15, bold=True, align=PP_ALIGN.CENTER)
    source(s, "기획력 항목은 아이디어의 크기가 아니라 '누가, 언제, 무엇을 하고, 언제 멈추는가' 로 채점된다.")

    # 13 역할·비용·위험 ─────────────────────────────────────────────────
    s = blank(prs); head(s, "역할 · 비용 · 위험", "농가가 하는 일은 두 가지뿐이다", 13)
    box(s, Inches(0.8), Inches(1.9), Inches(3.85), Inches(2.5), fill=CREAM, line=LINE)
    text(s, Inches(1.1), Inches(2.12), Inches(3.3), Inches(0.4), "역할", size=13, bold=True, color=ORANGE)
    for i, (who, what) in enumerate([("농가", "사진 촬영 · 수량 확인"), ("거점", "검수 · 등급 확정 · 포장"),
                                     ("천안시", "표준 가격 · 정산 · 판정"), ("개발팀", "시스템 운영 · 측정")]):
        y = Inches(2.6 + i * 0.44)
        text(s, Inches(1.1), y, Inches(1.0), Inches(0.38), who, size=13.5, bold=True)
        text(s, Inches(2.05), y, Inches(2.5), Inches(0.38), what, size=12.5, color=MUTED)
    box(s, Inches(4.9), Inches(1.9), Inches(3.85), Inches(2.5), fill=CREAM, line=LINE)
    text(s, Inches(5.2), Inches(2.12), Inches(3.3), Inches(0.4), "비용 (Phase 1)", size=13, bold=True, color=ORANGE)
    text(s, Inches(5.2), Inches(2.58), Inches(3.3), Inches(1.0),
         [[("신규 하드웨어", {"size": 13.5, "color": MUTED})],
          [("0원", {"size": 30, "bold": True, "color": ORANGE})]], spacing=1.2)
    text(s, Inches(5.2), Inches(3.5), Inches(3.35), Inches(0.8),
         "농가는 자기 휴대폰, 거점은 기존 시설, 모델은 서버 CPU 에서 돌아 GPU 가 필요 없습니다.",
         size=12, color=MUTED, spacing=1.25)
    box(s, Inches(9.0), Inches(1.9), Inches(3.65), Inches(2.5), fill=CREAM, line=LINE)
    text(s, Inches(9.3), Inches(2.12), Inches(3.1), Inches(0.4), "서버 운영", size=13, bold=True, color=ORANGE)
    text(s, Inches(9.3), Inches(2.58), Inches(3.1), Inches(1.0),
         [[("월", {"size": 13.5, "color": MUTED})], [("5–15만 원", {"size": 26, "bold": True})]], spacing=1.2)
    text(s, Inches(9.3), Inches(3.5), Inches(3.1), Inches(0.8),
         "소규모 인스턴스 1대. 사진은 외부 API 로 나가지 않습니다.", size=12, color=MUTED, spacing=1.25)

    text(s, Inches(0.8), Inches(4.62), Inches(11.9), Inches(0.4), "가장 큰 위험과 대응", size=15, bold=True)
    risks = [
        ("폰 사진에서 정확도 급락", "이미 예상하고 상한을 걸어둠 · 후보 3개 + 직접 선택으로 등록은 계속 가능"),
        ("등록은 되는데 안 팔림", "Go/No-Go 에 판매 성사율 50% 포함 · 미달 시 수요 확보 방식부터 재설계"),
        ("등급 분쟁", "AI 는 참고값, 확정은 거점 실물 검수 · 코드가 강제하고 화면에 표기"),
    ]
    for i, (r, a) in enumerate(risks):
        y = Inches(5.05 + i * 0.57)
        box(s, Inches(0.8), y, Inches(11.85), Inches(0.5), fill=RGBColor(0xF7, 0xF2, 0xE9), line=None)
        text(s, Inches(1.1), y + Inches(0.09), Inches(3.5), Inches(0.36), r, size=13, bold=True, color=RED)
        text(s, Inches(4.7), y + Inches(0.09), Inches(7.8), Inches(0.36), a, size=12.5, color=MUTED)
    source(s, "위험을 숨기면 질의응답에서 드러난다. 가장 큰 위험을 먼저 적고, 그 위험이 현실이 돼도 서비스가 멈추지 않게 설계했다.")

    # 14 마무리 ──────────────────────────────────────────────────────────
    s = blank(prs, INK)
    text(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.4),
         "정리", size=14, bold=True, color=ORANGE)
    text(s, Inches(0.9), Inches(1.65), Inches(11.5), Inches(1.6),
         [[("천안 농가 ", {"size": 30, "color": CREAM}),
           (f"{OUTSIDE:,}가구", {"size": 30, "bold": True, "color": ORANGE}),
           ("가 로컬푸드 체계 밖에 있습니다.", {"size": 30, "color": CREAM})],
          [("그들이 못 파는 이유는 판로가 아니라 등록입니다.", {"size": 30, "bold": True, "color": CREAM})]],
         spacing=1.3)
    facts = [
        ("작동하는 MVP", "사진 한 장 → 판매 시작", "오늘 시연 가능"),
        ("자동 테스트", f"{TESTS}개", "정책을 우회하면 실패한다"),
        ("재현 가능", "공개 데이터 + 스크립트", "숫자를 직접 다시 계산할 수 있다"),
    ]
    for i, (k, v, sub) in enumerate(facts):
        x = Inches(0.9 + i * 3.95)
        box(s, x, Inches(3.7), Inches(3.7), Inches(1.55), fill=RGBColor(0x2E, 0x2A, 0x25), line=None)
        text(s, x + Inches(0.3), Inches(3.9), Inches(3.1), Inches(0.35), k, size=12, color=RGBColor(0x9A, 0x92, 0x85))
        text(s, x + Inches(0.3), Inches(4.25), Inches(3.15), Inches(0.5), v, size=17, bold=True, color=CREAM)
        text(s, x + Inches(0.3), Inches(4.78), Inches(3.15), Inches(0.4), sub, size=11.5, color=RGBColor(0x9A, 0x92, 0x85))
    text(s, Inches(0.9), Inches(5.65), Inches(11.5), Inches(0.9),
         [[("요청드리는 것은 성환 거점 1곳과 농가 5호, 8주입니다.", {"size": 19, "color": CREAM})],
          [("그 8주가 끝나면 기준 미달 여부를 숫자로 보고드리고, 미달이면 확산하지 않겠습니다.",
            {"size": 16, "color": RGBColor(0xC9, 0xBF, 0xAE)})]], spacing=1.35)
    text(s, Inches(0.9), Inches(6.9), Inches(11.5), Inches(0.35),
         "ON-FARM · 2026년 천안시 AI·데이터 기반 정책 아이디어 경진대회 과제4",
         size=11, color=RGBColor(0x6A, 0x62, 0x56))

    path = OUT / "ON-FARM_천안_기획서.pptx"
    prs.save(path)
    return path


def to_pdf(pptx: Path) -> Path | None:
    soffice = Path(r"C:\Program Files\LibreOffice\program\soffice.exe")
    if not soffice.exists():
        print("LibreOffice 가 없어 PDF 변환을 건너뜁니다.")
        return None
    subprocess.run(
        [str(soffice), "--headless", "--convert-to", "pdf", "--outdir", str(OUT), str(pptx)],
        check=True, capture_output=True, timeout=300,
    )
    pdf = pptx.with_suffix(".pdf")
    return pdf if pdf.exists() else None


if __name__ == "__main__":
    deck = build()
    print(f"✔ {deck.relative_to(ROOT)}  ({deck.stat().st_size // 1024:,} KB)")
    print(f"  미참여율 {OUTSIDE_RATE:.1%} · 테스트 {TESTS}개 · 신뢰도 상한 {POLICY_CAP:.0%}")
    pdf = to_pdf(deck)
    if pdf:
        print(f"✔ {pdf.relative_to(ROOT)}  ({pdf.stat().st_size // 1024:,} KB)")
    sys.exit(0)
